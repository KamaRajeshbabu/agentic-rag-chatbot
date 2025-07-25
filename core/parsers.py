# core/parsers.py
from PyPDF2 import PdfReader
import pandas as pd
from docx import Document
from pptx import Presentation

def parse_pdf(file_path):
    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        print(f"[PDF Parser] Error: {e}")
    return text

def parse_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"[DOCX Parser] Error: {e}")
    return text

def parse_csv(file_path):
    text = ""
    try:
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
    except Exception as e:
        print(f"[CSV Parser] Error: {e}")
    return text

def parse_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[PPTX Parser] Error: {e}")
    return text

def parse_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"[TXT Parser] Error: {e}")
        return ""
